#!/usr/bin/env python3
"""
Generate PowerPoint presentation for "Generative AI for Beginners: Build Your First Agent with Azure AI Foundry"
Based on presentation-outline.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create the full presentation deck"""
    prs = Presentation()
    # 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define minimal color scheme
    AZURE_BLUE = RGBColor(0, 120, 212)
    DARK_GRAY = RGBColor(50, 50, 50)
    LIGHT_GRAY = RGBColor(150, 150, 150)
    WHITE = RGBColor(255, 255, 255)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "Generative AI for Beginners"
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = DARK_GRAY
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Build Your First Agent with Azure AI Foundry"
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = AZURE_BLUE
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(2), Inches(5), Inches(9.333), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.text = "From zero to AI hero - watch your first agent come alive in just a few lines of code"
    tagline_p.font.size = Pt(18)
    tagline_p.font.italic = True
    tagline_p.font.color.rgb = LIGHT_GRAY
    tagline_p.alignment = PP_ALIGN.CENTER

    # Your info placeholder
    info_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    info_frame = info_box.text_frame
    info_p = info_frame.paragraphs[0]
    info_p.text = "[Your Name] | [Your Title] | [Contact]"
    info_p.font.size = Pt(16)
    info_p.font.color.rgb = LIGHT_GRAY
    info_p.alignment = PP_ALIGN.CENTER

    # Slide 2: About Me
    slide2 = add_content_slide(prs, "About Me")
    add_bullet_points(slide2, [
        "[Brief intro - 30 seconds max]",
        "[Why you're excited about AI]",
        "[What makes you relatable to beginners]",
        "",
        "Keep it short and authentic!"
    ])

    # Slide 3: What We'll Build Today
    slide3 = add_content_slide(prs, "What We'll Build Today")
    add_bullet_points(slide3, [
        "✓ An AI agent in Azure AI Foundry",
        "✓ A working application connected to the agent",
        "✓ End-to-end demo you can use",
        "",
        "In the next hour: from nothing to working!",
        "",
        "You'll leave with code you can use"
    ])

    # Slide 4: What is Generative AI?
    slide4 = add_content_slide(prs, "What is Generative AI?")
    add_bullet_points(slide4, [
        "AI models that create new content",
        "",
        "Examples:",
        "  • Text (conversations, summaries, code)",
        "  • Images (art, designs)",
        "  • Audio (voice, music)",
        "",
        "Think: AI that can generate, not just analyze"
    ])

    # Slide 5: What is an AI Agent?
    slide5 = add_content_slide(prs, "What is an AI Agent?")
    add_bullet_points(slide5, [
        "AI that can take actions (not just chat)",
        "",
        "Agent vs. Simple Chatbot:",
        "  • Chatbot: answers questions",
        "  • Agent: performs tasks, makes decisions",
        "",
        "Real-world examples:",
        "  • Customer support that creates tickets",
        "  • Research assistant that finds information",
        "  • Code helper that writes and tests code"
    ])

    # Slide 6: Why Azure AI Foundry?
    slide6 = add_content_slide(prs, "Why Azure AI Foundry?")
    add_bullet_points(slide6, [
        "✓ All-in-one platform for building AI solutions",
        "",
        "✓ No need to be an ML expert",
        "",
        "✓ Enterprise-ready, secure, scalable",
        "",
        "✓ Built-in tools for testing and deployment",
        "",
        "Perfect for beginners and experts alike!"
    ])

    # Slide 7: Today's Architecture
    slide7 = add_content_slide(prs, "Today's Architecture")

    # Add simple text-based diagram
    diagram_box = slide7.shapes.add_textbox(Inches(3), Inches(2), Inches(7.333), Inches(3.5))
    diagram_frame = diagram_box.text_frame
    diagram_frame.word_wrap = True

    diagram_text = """
    USER
      ↓
    YOUR APPLICATION
      ↓
    AZURE AI FOUNDRY (Agent)
      ↓
    RESPONSE back to user

    "Don't worry, this is easier than it looks!"
    """

    diagram_p = diagram_frame.paragraphs[0]
    diagram_p.text = diagram_text.strip()
    diagram_p.font.size = Pt(28)
    diagram_p.font.color.rgb = DARK_GRAY
    diagram_p.alignment = PP_ALIGN.CENTER
    diagram_p.line_spacing = 1.5

    # Slide 8: Let's Build! Part 1
    slide8 = add_content_slide(prs, "Let's Build! Part 1 - The Agent")
    add_bullet_points(slide8, [
        "What we'll do:",
        "",
        "  ✓ Set up Azure AI Foundry",
        "  ✓ Create an agent",
        "  ✓ Configure prompts",
        "  ✓ Test in playground",
        "",
        "Minimal slides from here - let's code!"
    ])

    # Slide 9: Let's Build! Part 2
    slide9 = add_content_slide(prs, "Let's Build! Part 2 - The Application")
    add_bullet_points(slide9, [
        "What we'll do:",
        "",
        "  ✓ Create a simple front-end app",
        "  ✓ Connect to our agent",
        "  ✓ Handle user input/output",
        "  ✓ See it work end-to-end",
        "",
        "Keep watching - this is where it gets fun!"
    ])

    # Slide 10: What We Built Today
    slide10 = add_content_slide(prs, "What We Built Today")
    add_bullet_points(slide10, [
        "✓ Created an AI agent in Azure AI Foundry",
        "",
        "✓ Configured it with prompts",
        "",
        "✓ Connected it to a real application",
        "",
        "✓ Made it work end-to-end",
        "",
        "And you can do this too!"
    ])

    # Slide 11: Key Takeaways
    slide11 = add_content_slide(prs, "Key Takeaways")
    add_bullet_points(slide11, [
        "✓ Generative AI is accessible to everyone",
        "",
        "✓ Azure AI Foundry makes it easier",
        "",
        "✓ You don't need to be an AI expert",
        "",
        "✓ Start small, iterate, improve",
        "",
        "✓ The code is simpler than you think"
    ])

    # Slide 12: Next Steps
    slide12 = add_content_slide(prs, "Next Steps")
    add_bullet_points(slide12, [
        "Get the code: [GitHub link]",
        "",
        "Try Azure AI Foundry (free tier available)",
        "",
        "Learning resources:",
        "  • Microsoft Learn modules",
        "  • Azure AI Foundry documentation",
        "  • Community tutorials",
        "",
        "[Your contact for questions]"
    ])

    # Slide 13: Resources & Links
    slide13 = add_content_slide(prs, "Resources & Links")
    add_bullet_points(slide13, [
        "Azure AI Foundry:",
        "  ai.azure.com",
        "",
        "Documentation:",
        "  learn.microsoft.com/azure/ai-foundry",
        "",
        "Sample Code:",
        "  [Your GitHub repository]",
        "",
        "Contact:",
        "  [Your email/social media]"
    ])

    # Slide 14: Q&A
    slide14 = add_content_slide(prs, "Questions?")

    # Add large centered text
    qa_box = slide14.shapes.add_textbox(Inches(3.5), Inches(3), Inches(6.333), Inches(2))
    qa_frame = qa_box.text_frame
    qa_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    qa_p = qa_frame.paragraphs[0]
    qa_p.text = "Let's discuss!"
    qa_p.font.size = Pt(48)
    qa_p.font.color.rgb = DARK_GRAY
    qa_p.alignment = PP_ALIGN.CENTER

    # Slide 15: Thank You
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Thank you message
    thanks_box = slide15.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_p = thanks_frame.paragraphs[0]
    thanks_p.text = "Thank You!"
    thanks_p.font.size = Pt(72)
    thanks_p.font.bold = True
    thanks_p.font.color.rgb = DARK_GRAY
    thanks_p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide15.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.8))
    contact_frame = contact_box.text_frame
    contact_p = contact_frame.paragraphs[0]
    contact_p.text = "[Your Name] | [Email] | [LinkedIn/Twitter]"
    contact_p.font.size = Pt(24)
    contact_p.font.color.rgb = LIGHT_GRAY
    contact_p.alignment = PP_ALIGN.CENTER

    # Resources reminder
    resources_box = slide15.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9.333), Inches(0.6))
    resources_frame = resources_box.text_frame
    resources_p = resources_frame.paragraphs[0]
    resources_p.text = "Code & Resources: [GitHub/Website URL]"
    resources_p.font.size = Pt(20)
    resources_p.font.color.rgb = LIGHT_GRAY
    resources_p.alignment = PP_ALIGN.CENTER

    return prs


def add_content_slide(prs, title_text):
    """Add a slide with title and content area - minimal clean design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Minimal color scheme
    AZURE_BLUE = RGBColor(0, 120, 212)
    DARK_GRAY = RGBColor(50, 50, 50)

    # Title - simple, no background
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = AZURE_BLUE
    title_p.alignment = PP_ALIGN.LEFT

    return slide


def add_bullet_points(slide, bullet_points):
    """Add bullet points to a slide - minimal clean design"""
    DARK_GRAY = RGBColor(50, 50, 50)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.font.size = Pt(26)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

        # Set bullet level based on indentation
        if point.startswith("  •") or point.startswith("  "):
            p.level = 1
        else:
            p.level = 0


def main():
    """Main function to create and save presentation"""
    print("Creating PowerPoint presentation...")
    prs = create_presentation()

    output_file = "AI_Foundry_Beginners_Presentation_Widescreen.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Format: 16:9 widescreen")
    print(f"✓ Design: Clean, minimal, no background colors")
    print("\nNext steps:")
    print("1. Open the presentation in PowerPoint")
    print("2. Replace placeholder text with your information")
    print("3. Add images/diagrams as needed")
    print("4. Customize as needed")


if __name__ == "__main__":
    main()
